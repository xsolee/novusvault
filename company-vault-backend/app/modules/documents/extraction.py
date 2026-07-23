"""Mime-type dispatch for native text extraction. OCR itself lives behind
providers/ocr — this module only extracts what's natively available and
flags when a document is too sparse (scanned/low-text) to skip OCR."""

import csv
import io

import fitz  # PyMuPDF
import openpyxl
from docx import Document as DocxDocument
from pptx import Presentation

_IMAGE_MIME_TYPES = {"image/png", "image/jpeg", "image/jpg", "image/tiff", "image/bmp", "image/webp"}
_LOW_TEXT_CHARS_PER_PAGE_THRESHOLD = 100


def extract_native_text(file_bytes: bytes, mime_type: str) -> str:
    if mime_type == "application/pdf":
        return _extract_pdf_text(file_bytes)
    if mime_type in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ):
        return _extract_docx_text(file_bytes)
    if mime_type in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ):
        return _extract_xlsx_text(file_bytes)
    if mime_type in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ):
        return _extract_pptx_text(file_bytes)
    if mime_type == "text/csv":
        return _extract_csv_text(file_bytes)
    if mime_type.startswith("text/"):
        return file_bytes.decode("utf-8", errors="replace")
    if mime_type in _IMAGE_MIME_TYPES:
        return ""  # no native text — OCR stage handles this entirely
    return ""


def is_image(mime_type: str) -> bool:
    return mime_type in _IMAGE_MIME_TYPES


def is_low_text_pdf(file_bytes: bytes, extracted_text: str) -> bool:
    if not file_bytes:
        return False
    page_count = get_pdf_page_count(file_bytes)
    if page_count == 0:
        return False
    return (len(extracted_text) / page_count) < _LOW_TEXT_CHARS_PER_PAGE_THRESHOLD


def get_pdf_page_count(file_bytes: bytes) -> int:
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        return doc.page_count


def rasterize_pdf_pages(file_bytes: bytes, *, dpi: int = 150) -> list[bytes]:
    """Renders each PDF page to PNG bytes for OCR fallback on scanned/low-text PDFs."""
    pages: list[bytes] = []
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for page in doc:
            pixmap = page.get_pixmap(dpi=dpi)
            pages.append(pixmap.tobytes("png"))
    return pages


def _extract_pdf_text(file_bytes: bytes) -> str:
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        return "\n\n".join(page.get_text() for page in doc)


def _extract_docx_text(file_bytes: bytes) -> str:
    doc = DocxDocument(io.BytesIO(file_bytes))
    return "\n".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text)


def _extract_xlsx_text(file_bytes: bytes) -> str:
    workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True, read_only=True)
    lines: list[str] = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell) for cell in row if cell is not None]
            if cells:
                lines.append(" | ".join(cells))
    return "\n".join(lines)


def _extract_pptx_text(file_bytes: bytes) -> str:
    presentation = Presentation(io.BytesIO(file_bytes))
    lines: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.append(shape.text_frame.text)
    return "\n".join(lines)


def _extract_csv_text(file_bytes: bytes) -> str:
    text = file_bytes.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    return "\n".join(" | ".join(row) for row in reader)
